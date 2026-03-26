import frappe
from pathlib import Path


class FileTranscriber:

	def __init__(self):
		self.logger = frappe.logger("os_lmsa", allow_site=True)

	def transcribe(self, file_path: str, file_type: str) -> str:
		"""
		Extract text content from a file.

		Args:
			file_path: Path to the file (absolute or relative to Frappe site).
			file_type: File type/extension (e.g. "pdf", "docx", "txt").

		Returns:
			Extracted text content as a string.
		"""
		file_type = file_type.lower().strip().lstrip(".")

		handlers = {
			"pdf": self._transcribe_pdf,
			"docx": self._transcribe_docx,
			"txt": self._transcribe_txt,
			"pptx": self._transcribe_pptx,
		}

		handler = handlers.get(file_type)
		if handler is None:
			self.logger.warning(f"Unsupported file type: {file_type}")
			return ""

		if file_path.startswith("/private/"):
			resolved_path = Path(frappe.get_site_path(file_path.lstrip("/"))).resolve()
		else:
			resolved_path = Path(frappe.get_site_path("public", file_path.lstrip("/"))).resolve()
		
		if not resolved_path.exists():
			self.logger.warning(f"File not found: {resolved_path}")
			return ""

		self.logger.info(f"Transcribing file {resolved_path} (type={file_type})")
		text = handler(resolved_path)
		self.logger.info(f"Transcribed {len(text)} chars from {resolved_path}")
		return text

	

	def _transcribe_pdf(self, file_path: str) -> str:
		import fitz  # PyMuPDF

		doc = fitz.open(file_path)
		parts = []
		for page in doc:
			text = page.get_text()
			if text.strip():
				parts.append(text.strip())
		doc.close()
		return "\n\n".join(parts)

	def _transcribe_docx(self, file_path: str) -> str:
		import docx

		doc = docx.Document(file_path)
		parts = []
		for para in doc.paragraphs:
			if para.text.strip():
				parts.append(para.text.strip())
		return "\n\n".join(parts)

	def _transcribe_txt(self, file_path: str) -> str:
		return Path(file_path).read_text(encoding="utf-8")

	def _transcribe_pptx(self, file_path: str) -> str:
		from pptx import Presentation

		prs = Presentation(file_path)
		parts = []
		for slide in prs.slides:
			for shape in slide.shapes:
				if shape.has_text_frame:
					text = shape.text_frame.text.strip()
					if text:
						parts.append(text)
		return "\n\n".join(parts)
